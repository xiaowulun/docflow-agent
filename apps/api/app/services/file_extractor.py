"""
文件内容提取工具

支持从各种文档格式中提取文本内容，供 LLM 使用。
"""

from pathlib import Path


def extract_text(file_path: Path) -> str:
    """从文件中提取文本内容"""
    suffix = file_path.suffix.lower()

    extractors = {
        ".txt": _extract_txt,
        ".md": _extract_txt,
        ".csv": _extract_txt,
        ".json": _extract_txt,
        ".xml": _extract_txt,
        ".html": _extract_html,
        ".htm": _extract_html,
        ".docx": _extract_docx,
        ".doc": _extract_doc,
        ".xlsx": _extract_xlsx,
        ".xls": _extract_xlsx,
        ".pdf": _extract_pdf,
        ".pptx": _extract_pptx,
        ".ppt": _extract_pptx,
    }

    extractor = extractors.get(suffix)
    if not extractor:
        return f"[无法提取 {suffix} 文件内容]"

    try:
        return extractor(file_path)
    except Exception as e:
        return f"[提取 {suffix} 文件内容失败: {e}]"


def _extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_html(path: Path) -> str:
    """简单提取 HTML 文本（去掉标签）"""
    import re
    text = path.read_text(encoding="utf-8", errors="replace")
    text = re.sub(r"<script[^>]*>.*?</script>", "", text, flags=re.DOTALL)
    text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _extract_doc(path: Path) -> str:
    """提取旧版 .doc 文件（OLE2 格式）"""
    import subprocess
    import tempfile
    
    # 使用 macOS 的 textutil 工具将 .doc 转换为 .txt
    try:
        with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as tmp:
            tmp_path = tmp.name
        
        # textutil 是 macOS 自带的文档转换工具
        result = subprocess.run(
            ['textutil', '-convert', 'txt', '-output', tmp_path, str(path)],
            capture_output=True,
            text=True,
            timeout=30
        )
        
        if result.returncode == 0:
            # 读取转换后的文本文件
            converted_text = Path(tmp_path).read_text(encoding='utf-8', errors='ignore')
            Path(tmp_path).unlink()  # 删除临时文件
            return converted_text.strip()
        else:
            Path(tmp_path).unlink(missing_ok=True)
            return f"[textutil 转换失败: {result.stderr}]"
    
    except FileNotFoundError:
        # textutil 不存在（非 macOS 系统）
        return "[系统不支持 .doc 文件提取，请安装 antiword 或转换为 .docx 格式]"
    except subprocess.TimeoutExpired:
        return "[.doc 文件处理超时]"
    except Exception as e:
        return f"[提取 .doc 文件失败: {e}]"


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"=== 工作表: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            lines.append("\t".join(cells))
    wb.close()
    return "\n".join(lines)


def _extract_pdf(path: Path) -> str:
    import pdfplumber
    lines = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text:
                lines.append(f"--- 第 {i + 1} 页 ---")
                lines.append(text)
    return "\n".join(lines)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    lines = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"--- 幻灯片 {i + 1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


# 文件类型分类
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
DOCUMENT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm",
    ".docx", ".doc", ".xlsx", ".xls", ".pdf", ".pptx", ".ppt",
}
ALL_ALLOWED = IMAGE_EXTENSIONS | DOCUMENT_EXTENSIONS
